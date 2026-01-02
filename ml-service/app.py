from flask import Flask, request, jsonify
from flask_cors import CORS
from pypdf import PdfReader
from openai import OpenAI
import os
import tempfile
import pandas as pd
import joblib

# Initialize Flask App
app = Flask(__name__)
CORS(app)

# Configuration
# TODO: Move this API key to an environment variable for security
# e.g., os.environ.get("OPENAI_API_KEY")
OPENAI_API_KEY = "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"
client = OpenAI(api_key=OPENAI_API_KEY)

# Load ML Model
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
MODEL_PATH = os.path.join(BASE_DIR, "performance_model.pkl")
model = None

if os.path.exists(MODEL_PATH):
    try:
        model = joblib.load(MODEL_PATH)
        print("✅ ML Model loaded successfully")
    except Exception as e:
        print(f"⚠️ Error loading model: {e}")
else:
    print("⚠️ Model file not found. Prediction features will not work.")

# --- Helper Functions ---

def pdf_to_text(pdf_file):
    """Extract text from a PDF file."""
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def ppt_to_text(ppt_file):
    """Extract text from a PPT/PPTX file."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx' if ppt_file.filename.endswith('.pptx') else '.ppt') as tmp_file:
            ppt_file.save(tmp_file.name)
            tmp_path = tmp_file.name
        
        try:
            from pptx import Presentation
            prs = Presentation(tmp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            os.unlink(tmp_path)
            return text
        except ImportError:
            os.unlink(tmp_path)
            return None
    except Exception as e:
        print(f"Error extracting text from PPT: {e}")
        return None

# --- Routes ---

@app.route("/", methods=["GET"])
def home():
    return jsonify({"message": "EduLens ML Service is running!"})

@app.route("/extract_text", methods=["POST"])
def extract_text():
    if "file" not in request.files:
        return jsonify({"error": "Please upload a file"}), 400

    file = request.files["file"]
    filename = file.filename.lower()

    try:
        if filename.endswith('.pdf'):
            text = pdf_to_text(file)
        elif filename.endswith('.ppt') or filename.endswith('.pptx'):
            text = ppt_to_text(file)
            if text is None:
                return jsonify({
                    "error": "PPT extraction requires python-pptx library. Please install it with: pip install python-pptx"
                }), 400
        else:
            return jsonify({"error": "Unsupported file type. Please upload PDF or PPT/PPTX"}), 400

        if not text.strip():
            return jsonify({"error": "No readable text found in file"}), 400

        return jsonify({"text": text})
    except Exception as e:
        print(f"Error in extract_text: {str(e)}")
        return jsonify({"error": f"Error extracting text: {str(e)}"}), 500


@app.route("/summarize_pdf", methods=["POST"])
def summarize_pdf():
    if "file" not in request.files:
        return jsonify({"error": "Please upload a PDF file"}), 400

    pdf = request.files["file"]
    text = pdf_to_text(pdf)

    if not text.strip():
        return jsonify({"error": "No readable text in PDF"}), 400

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Summarize this text:\n\n{text}"}
        ]
    )

    summary = response.choices[0].message.content
    return jsonify({"summary": summary})


@app.route("/quiz_from_pdf", methods=["POST"])
def quiz_pdf():
    if "file" not in request.files:
        return jsonify({"error": "Please upload a PDF file"}), 400

    pdf = request.files["file"]
    text = pdf_to_text(pdf)

    if not text.strip():
        return jsonify({"error": "No readable text in PDF"}), 400

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Generate 10-20 single-choice questions (MCQ) for a university exam based on the text below. "
        f"Strictly follow this format for every question:\n\n"
        f"1. Question text here?\n"
        f"A) Option 1\n"
        f"B) Option 2\n"
        f"C) Option 3\n"
        f"D) Option 4\n"
        f"Answer:-: A) Option 1\n\n"
        f"Text to use:\n{text}"}
        ]
    )

    summary = response.choices[0].message.content
    return jsonify({"quiz": summary})


@app.route("/generate_flashcards", methods=["POST"])
def generate_flashcards():
    data = request.get_json()
    text = data.get("text", "")
    
    if not text.strip():
        return jsonify({"error": "Text content is required"}), 400
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Generate 10-15 flashcards from the following text. "
             f"Format each flashcard as:\n"
             f"Front: [Question or term]\n"
             f"Back: [Answer or definition]\n\n"
             f"Text:\n{text}"}
        ]
    )
    
    flashcards_text = response.choices[0].message.content
    flashcards = []
    lines = flashcards_text.split('\n')
    current_card = {}
    
    for line in lines:
        line = line.strip()
        if line.startswith('Front:'):
            if current_card:
                flashcards.append(current_card)
            current_card = {'front': line.replace('Front:', '').strip(), 'back': ''}
        elif line.startswith('Back:'):
            if current_card:
                current_card['back'] = line.replace('Back:', '').strip()
    
    if current_card:
        flashcards.append(current_card)
    
    return jsonify({"flashcards": flashcards})


@app.route("/explain_concept", methods=["POST"])
def explain_concept():
    data = request.get_json()
    concept = data.get("concept", "")
    level = data.get("level", "beginner")
    
    if not concept.strip():
        return jsonify({"error": "Concept is required"}), 400
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Explain the concept '{concept}' in a {level} level. "
             f"Make it clear, concise, and easy to understand. Include examples if helpful."}
        ]
    )
    
    explanation = response.choices[0].message.content
    return jsonify({"explanation": explanation})


@app.route("/generate_study_plan", methods=["POST"])
def generate_study_plan():
    data = request.get_json()
    subjects = data.get("subjects", [])
    days = data.get("days", 7)
    hours_per_day = data.get("hours_per_day", 2)
    
    if not subjects:
        return jsonify({"error": "At least one subject is required"}), 400
    
    subjects_str = ", ".join(subjects)
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Create a {days}-day study plan for the following subjects: {subjects_str}. "
             f"Available study time: {hours_per_day} hours per day. "
             f"Format as a daily schedule with specific topics and activities for each day. "
             f"Make it realistic and balanced."}
        ]
    )
    
    study_plan = response.choices[0].message.content
    return jsonify({"study_plan": study_plan})


@app.route("/generate_notes", methods=["POST"])
def generate_notes():
    data = request.get_json()
    topic = data.get("topic", "")
    detail_level = data.get("detail_level", "medium")
    
    if not topic.strip():
        return jsonify({"error": "Topic is required"}), 400
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Create comprehensive study notes on the topic: '{topic}'. "
             f"Detail level: {detail_level}. "
             f"Format with clear headings, bullet points, and key concepts highlighted. "
             f"Make it suitable for exam preparation."}
        ]
    )
    
    notes = response.choices[0].message.content
    return jsonify({"notes": notes})


@app.route("/explain_answer", methods=["POST"])
def explain_answer():
    data = request.get_json()
    question = data.get("question", "")
    answer = data.get("answer", "")
    is_correct = data.get("is_correct", True)
    
    if not question or not answer:
        return jsonify({"error": "Question and answer are required"}), 400
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Question: {question}\n"
             f"Answer: {answer}\n"
             f"Provide a clear explanation of why this answer is {'correct' if is_correct else 'incorrect'}. "
             f"Explain the concept and provide additional context if helpful."}
        ]
    )
    
    explanation = response.choices[0].message.content
    return jsonify({"explanation": explanation})


@app.route("/ai_chat", methods=["POST"])
def ai_chat():
    data = request.get_json()
    message = data.get("message", "")
    context = data.get("context", "")
    
    if not message.strip():
        return jsonify({"error": "Message is required"}), 400
    
    system_prompt = "You are a helpful AI study assistant. Answer questions clearly and concisely. Help students understand concepts, solve problems, and study effectively."
    
    user_content = message
    if context:
        user_content = f"Context: {context}\n\nQuestion: {message}"
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_content}
        ]
    )
    
    reply = response.choices[0].message.content
    return jsonify({"reply": reply})


@app.route("/analyze_difficulty", methods=["POST"])
def analyze_difficulty():
    data = request.get_json()
    questions = data.get("questions", [])
    
    if not questions:
        return jsonify({"error": "Questions are required"}), 400
    
    questions_text = "\n".join([f"{i+1}. {q}" for i, q in enumerate(questions)])
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "user", "content": f"Analyze the difficulty level of these questions and categorize each as 'easy', 'medium', or 'hard'. "
             f"Provide a brief explanation for each:\n\n{questions_text}"}
        ]
    )
    
    analysis = response.choices[0].message.content
    return jsonify({"analysis": analysis})

# --- Merged from route-ml.py (Performance Prediction) ---
@app.route('/predict_performance', methods=['POST'])
def predict():
    if model is None:
         return jsonify({"error": "Model not loaded. Cannot predict."}), 500
         
    try:
        data = request.get_json()
        df = pd.DataFrame([data])
        # Note: Ensure the incoming data matches the model's expected features
        # 'attendance_pct','assignments_avg','midterm','final'
        prediction = model.predict(df)[0]
        return jsonify({"predicted_performance": int(prediction)})
    except Exception as e:
        print(f"Prediction Error: {e}")
        return jsonify({"error": f"Prediction failed: {str(e)}"}), 500


if __name__ == "__main__":
    print("Starting EduLens ML Service on port 5000...")
    app.run(port=5000)
